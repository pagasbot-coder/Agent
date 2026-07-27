#!/usr/bin/env python3
"""Travel+ Natural interview deck v1.14.1 — design polish (UI/UX + Copy)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os
import shutil
from datetime import datetime, timezone

W = Inches(13.333)
H = Inches(7.5)

BG = RGBColor(0xF5, 0xF3, 0xEF)
ACCENT = RGBColor(0x2A, 0x55, 0x4A)
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x5C, 0x5C, 0x5C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xD4, 0xCE, 0xC4)

MARGIN_X = Inches(0.5)
CONTENT_TOP = Inches(1.05)
CARD_H = Inches(5.55)
GAP = Inches(0.28)
CARD_W = (W - 2 * MARGIN_X - GAP) / 2
PAD = Inches(0.32)


def set_run(run, size=16, bold=False, color=INK, name="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    rFonts = rPr.find(qn("a:rFonts"))
    if rFonts is None:
        rFonts = etree.SubElement(rPr, qn("a:rFonts"))
    rFonts.set("ascii", name)
    rFonts.set("hAnsi", name)
    rFonts.set("cs", name)


def add_bg(slide):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    shp.fill.solid()
    shp.fill.fore_color.rgb = BG
    shp.line.fill.background()
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.1), H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def add_footer(slide, page, total=12):
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN_X, Inches(6.85), W - 2 * MARGIN_X, Inches(0.015)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()
    box = slide.shapes.add_textbox(MARGIN_X, Inches(6.95), Inches(10.5), Inches(0.35))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Travel+ Natural  ·  пульт v1.14  ·  собеседование  ·  не оферта"
    set_run(r, 11, False, MUTED)
    box2 = slide.shapes.add_textbox(Inches(11.5), Inches(6.95), Inches(1.4), Inches(0.35))
    tf2 = box2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run()
    r2.text = f"{page} / {total}"
    set_run(r2, 11, False, MUTED)


def add_title(slide, text):
    box = slide.shapes.add_textbox(MARGIN_X, Inches(0.28), Inches(12.3), Inches(0.55))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    set_run(r, 26, True, ACCENT)
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN_X, Inches(0.88), Inches(2.2), Inches(0.04)
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = ACCENT
    rule.line.fill.background()


def add_card(slide, left, top, width, height):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = CARD
    shp.line.color.rgb = LINE
    shp.line.width = Pt(1)
    try:
        shp.adjustments[0] = 0.06
    except Exception:
        pass
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.08))
    strip.fill.solid()
    strip.fill.fore_color.rgb = ACCENT
    strip.line.fill.background()
    return shp


def write_block(slide, left, top, width, height, lines):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(item.get("before", 0))
        p.space_after = Pt(item.get("after", 8))
        p.line_spacing = 1.15
        r = p.add_run()
        r.text = item["text"]
        set_run(r, item.get("size", 16), item.get("bold", False), item.get("color", INK))
    return box


def card_header(text):
    return {"text": text, "size": 15, "bold": True, "color": ACCENT, "after": 14}


def body(text, **kw):
    d = {"text": text, "size": 16, "bold": False, "color": INK, "after": 8}
    d.update(kw)
    return d


def muted(text, **kw):
    d = {"text": text, "size": 14, "bold": False, "color": MUTED, "after": 6}
    d.update(kw)
    return d


def strong(text, **kw):
    d = {"text": text, "size": 16, "bold": True, "color": INK, "after": 8}
    d.update(kw)
    return d


def two_cards(slide):
    left = MARGIN_X
    right = MARGIN_X + CARD_W + GAP
    add_card(slide, left, CONTENT_TOP, CARD_W, CARD_H)
    add_card(slide, right, CONTENT_TOP, CARD_W, CARD_H)
    tx = left + PAD
    ty = CONTENT_TOP + Inches(0.28)
    tw = CARD_W - 2 * PAD
    th = CARD_H - Inches(0.45)
    return (tx, ty, tw, th), (right + PAD, ty, tw, th)


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]

    # 1 Title
    s = prs.slides.add_slide(blank)
    add_bg(s)
    panel = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(5.1), H)
    panel.fill.solid()
    panel.fill.fore_color.rgb = ACCENT
    panel.line.fill.background()
    write_block(
        s,
        Inches(0.55),
        Inches(2.15),
        Inches(4.2),
        Inches(3.2),
        [
            {"text": "Павел Биджиев", "size": 30, "bold": True, "color": WHITE, "after": 14},
            {"text": "Продакт / бренд-логика", "size": 17, "color": WHITE, "after": 8},
            {
                "text": "Физические товары · HoReCa amenities",
                "size": 14,
                "color": RGBColor(0xC8, 0xDD, 0xD6),
                "after": 0,
            },
        ],
    )
    write_block(
        s,
        Inches(5.7),
        Inches(2.0),
        Inches(6.9),
        Inches(3.8),
        [
            {"text": "Travel+ Natural", "size": 32, "bold": True, "color": ACCENT, "after": 12},
            {"text": "Своя натуральная линейка:", "size": 19, "color": INK, "after": 4},
            {"text": "от рамки до пилота", "size": 19, "color": INK, "after": 18},
            {
                "text": "Подготовка · АК-Сервис / Travel+ · Кириши",
                "size": 14,
                "color": MUTED,
                "after": 6,
            },
            {"text": "Пульт бренда v1.14 · не оферта", "size": 14, "color": MUTED, "after": 0},
        ],
    )

    # 2 Role
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Как понял роль")
    add_footer(s, 2)
    L, R = two_cards(s)
    write_block(
        s,
        *L,
        [
            card_header("Travel+"),
            body("Производство и one-stop на номер:"),
            body("• Косметика"),
            body("• Тапочки"),
            body("• Наборы", after=14),
            strong("Не IT-roadmap."),
            body("Товар от идеи до полки и маржи."),
        ],
    )
    write_block(
        s,
        *R,
        [
            card_header("Роль здесь"),
            body("• Ассортимент и новинки"),
            body("• Позиция линейки в доме брендов"),
            body("• Продукт и доказательства до витрины"),
            body("• Связка: технолог → продажи → экономика → поставка"),
        ],
    )

    # 3 Gap
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Дыра на витрине")
    add_footer(s, 3)
    L, R = two_cards(s)
    write_block(
        s,
        *L,
        [
            card_header("Рынок"),
            body("Конкуренты: ЕТС Natural · MEZO · NS"),
            muted("(NS — чужой premium в каталоге)", after=12),
            body("Travel+ силён в масс и среднем."),
            strong("Слабое место: своя Natural с документами."),
        ],
    )
    write_block(
        s,
        *R,
        [
            card_header("Ориентир цен (не оферта)"),
            body("Hotel Line ≈ 12–15 ₽ / 30 мл"),
            body("ЕТС Natural ≈ 18 ₽ / 25 мл"),
            body("Natura Siberica — выше", after=14),
            strong("Ставка: между HL и NS"),
            body("Натуральность · документы · один поставщик"),
        ],
    )

    # 4 Portfolio
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Дом линеек")
    add_footer(s, 4)
    L, R = two_cards(s)
    write_block(
        s,
        *L,
        [
            card_header("Карта"),
            body("Hotel Line → цена"),
            body("Fleur / Aquatique / La Nuit → дизайн и аромат"),
            strong("Natural → натуральность + документы"),
            body("Natura Siberica → чужой premium"),
            body("СТМ (лого отеля) → отдельный проект"),
        ],
    )
    write_block(
        s,
        *R,
        [
            card_header("Мостик для продаж"),
            body("Дизайн и аромат — Fleur."),
            body("Натуральность и документы — Natural."),
            strong("Не конкурируют."),
            body("Могут жить на разных этажах одного объекта."),
        ],
    )

    # 5 Promise
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Обещание + StoryBrand")
    add_footer(s, 5)
    L, R = two_cards(s)
    write_block(
        s,
        *L,
        [
            card_header("Обещание"),
            strong("Мягкий уход с понятной натуральностью"),
            body("для отеля, которому важны и гость, и документы.", after=12),
            muted("Доказательства:", after=6),
            body("• Состав в границах"),
            body("• Пакет документов Travel+"),
            body("• Ощущение в номере"),
            body("• Один поставщик на номер"),
            body("• Блок «простыми словами» на этикетке"),
        ],
    )
    write_block(
        s,
        *R,
        [
            card_header("BrandScript"),
            body("Герой: закупщик eco / wellness"),
            body("Проводник: Travel+ Natural"),
            body("План: сегмент → образец + документы → пилот", after=14),
            muted("Гость после дороги:", after=6),
            strong("Мягче, чем ждал."),
            body("Понятная натуральность без громких обещаний."),
        ],
    )

    # 6 DoD
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Границы до рынка — DoD этапа 2")
    add_footer(s, 6)
    L, R = two_cards(s)
    write_block(
        s,
        *L,
        [
            card_header("Пока красное — витрину не открываем"),
            body("□ Метод и % натуральности"),
            body("□ Пакет как у HL / Fleur + лист Natural"),
            body("□ Origin зафиксирован"),
            body("□ Слепой тест запаха с гостями"),
            body("□ Этикетка «простыми словами»"),
            body("□ Себестоимость и коридор цены"),
            body("□ Срок поставки / мин. заказ / брак"),
        ],
    )
    write_block(
        s,
        *R,
        [
            card_header("Правила"),
            body("На старте без COSMOS, если знака нет — честно."),
            strong("Без DoD — не КП и не live GTM.", after=12),
            body("Этапы 3–6 — учебный сценарий,"),
            body("не факт завода."),
        ],
    )

    # 7 Pitch
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Питч закупщику (30 секунд)")
    add_footer(s, 7)
    add_card(s, MARGIN_X, CONTENT_TOP, W - 2 * MARGIN_X, Inches(4.6))
    write_block(
        s,
        MARGIN_X + PAD,
        CONTENT_TOP + Inches(0.35),
        W - 2 * MARGIN_X - 2 * PAD,
        Inches(4.0),
        [
            body(
                "«Вам нужно выглядеть зеленее — и спокойно пройти и гостя, и документы.",
                size=17,
                after=10,
            ),
            body("Travel+ Natural: пакет документов как у Hotel Line и Fleur", size=17, after=4),
            muted("(декларация / СГР, INCI, протоколы по запросу),", after=8),
            body(
                "прозрачный состав с блоком «простыми словами», один договор на номер.",
                size=17,
                after=10,
            ),
            strong("Три шага: eco / wellness → образец → пилот.»", size=17),
        ],
    )
    write_block(
        s,
        MARGIN_X,
        Inches(5.9),
        W - 2 * MARGIN_X,
        Inches(0.6),
        [muted("Отрицания (COSMOS, origin, цена Hotel Line) — только в возражениях.")],
    )

    # 8 Objections
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Поле: возражения и риски")
    add_footer(s, 8)
    L, R = two_cards(s)
    write_block(
        s,
        *L,
        [
            card_header("Возражения"),
            body("COSMOS обязателен → Natural в это КП не ставим"),
            body("«Natural по цене Hotel Line» → нет"),
            body("СТМ / лого отеля → отдельный проект"),
            body("«Где произведено?» → только зафиксированный origin"),
        ],
    )
    write_block(
        s,
        *R,
        [
            card_header("Полевая дисциплина"),
            body("Реестр рисков — один лист"),
            muted("(14 пунктов в пульте)", after=12),
            body("Сегмент не размываем."),
            strong("Не продаём Natural в эконом Hotel Line."),
        ],
    )

    # 9 Pilot
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Пилот: успех, провал, выход")
    add_footer(s, 9)
    L, R = two_cards(s)
    write_block(
        s,
        *L,
        [
            card_header("Успех (учебный порог)"),
            body("• 3 пилотных клиента"),
            body("• Повторный заказ у ≥ 2 из 3"),
            body("• Мягкость и ощущение натуральности ≥ 70%", after=10),
            strong("• Natural не ушёл в эконом Hotel Line", after=14),
            muted("Цифры сценария — учебный пример,"),
            muted("не отчёт завода."),
        ],
    )
    write_block(
        s,
        *R,
        [
            card_header("Провал и выход"),
            body("Повтор < 2 из 3 или метрики < 60%"),
            strong("→ Останавливаем масштаб", after=12),
            body("Диагноз бренд-менеджера за 2 недели"),
            body("Второй круг снова провален"),
            strong("→ Закрываем линейку"),
        ],
    )

    # 10 Logistics
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Поставки и экономика (Китай)")
    add_footer(s, 10)
    L, R = two_cards(s)
    write_block(
        s,
        *L,
        [
            card_header("Экономика"),
            body("Прайс — после себестоимости на складе."),
            body("Ориентир: выше Hotel Line · рядом с ЕТС · ниже NS.", after=12),
            strong("Не воюем копейкой."),
            body("Выигрываем пакетом и ощущением."),
        ],
    )
    write_block(
        s,
        *R,
        [
            card_header("Если Китай или гибрид"),
            body("Сравниваем заводы в FOB"),
            body("Логистику до Кириши считаем отдельно"),
            body("Полная цена = FOB + фрахт + таможня + довоз", after=10),
            strong("Срок пилота — «на складе»"),
            body("Origin не ясен → не говорим «российское»"),
        ],
    )

    # 11 What's in pult
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Что уже есть в пульте")
    add_footer(s, 11)
    L, R = two_cards(s)
    write_block(
        s,
        *L,
        [
            card_header("Содержание"),
            body("Этапы 0–2 — разведка и правила"),
            body("Этапы 3–6 — сценарий «как вести дело»"),
            body("StoryBrand · DoD · риски · возврат"),
            body("Мостик портфеля Travel+"),
        ],
    )
    write_block(
        s,
        *R,
        [
            card_header("Справочники"),
            body("Глоссарий бренда + глоссарий ВЭД"),
            body("Раздел L: путь · Incoterms · RFQ", after=14),
            body("Статус: рамка для поля и собеседования."),
            strong("v2.0 — когда технолог даст цифры."),
        ],
    )

    # 12 Close
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Итог и следующий шаг")
    add_footer(s, 12)
    L, R = two_cards(s)
    write_block(
        s,
        *L,
        [
            card_header("Собрал"),
            body("ДНК + StoryBrand + DoD"),
            body("Полевая защита"),
            body("Чеклист поставок при Китае (раздел L)", after=14),
            muted("Не прайс завода. Не КП."),
        ],
    )
    write_block(
        s,
        *R,
        [
            card_header("Готов обсудить"),
            body("1) Приоритет блокеров этапа 2 с технологом"),
            body("2) Пилотный сегмент eco / wellness"),
            body("3) Как вести продажи, чтобы Natural не ушёл в эконом Hotel Line", after=14),
            strong("Спасибо. Вопросы."),
        ],
    )

    out = "knowledge-base/presentation-travelplus-interview.pptx"
    prs.save(out)
    dests = [
        "АК/presentation-travelplus-interview.pptx",
        "/opt/cursor/artifacts/presentation-travelplus-interview.pptx",
        "/opt/cursor/artifacts/АК/presentation-travelplus-interview.pptx",
        "/home/ubuntu/АК/presentation-travelplus-interview.pptx",
        "/home/ubuntu/Desktop/АК/presentation-travelplus-interview.pptx",
    ]
    for dest in dests:
        os.makedirs(os.path.dirname(dest) or ".", exist_ok=True)
        shutil.copy2(out, dest)
    print("saved", out, "slides", len(prs.slides))
    return out


if __name__ == "__main__":
    build()
